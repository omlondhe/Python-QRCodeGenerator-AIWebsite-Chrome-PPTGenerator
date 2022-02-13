# import Presentation class
from pptx import Presentation

# Creating presentation object
root = Presentation()

# Creating slide layout
first_slide_layout = root.slide_layouts[0]

# Creating slide object to add slide created above
slide = root.slides.add_slide(first_slide_layout)

# Adding title and subtitle
slide.shapes.title.text = "Hello World"
slide.placeholders[1].text = "This PPT is created using Python"

# Saving file
root.save("Output.pptx")

print("Created")
